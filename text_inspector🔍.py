import streamlit as st #build interactive web app
import pdfplumber# to read pdf file via upload by user
from docx import Document#to read word or .doc file via upload by user
from pptx import Presentation#to read pptx file via upload by user
import re# pattern matching in text
from collections import Counter# how many times each item appears in a list
import io#Handles files in memory without saving to disk.

### Let's start our text_inspector🔍

st.set_page_config(
    page_title="Text_Inspector",
    page_icon="🔍"
)

st.markdown("<h1 style='text-align: center;'>Text Inspector 🔍</h1>", unsafe_allow_html=True)

st.markdown("""
    <style>
    .stApp {
        background-color: #C1E1C1;
            
    }
    </style>
""", unsafe_allow_html=True)

st.markdown("welcome to text inspector🔍(●'◡'●)")

uploaded_file = st.file_uploader("**Upload your file**", type=["txt", "pdf", "docx", "pptx"])
if uploaded_file is not None:
    st.success("**File uploaded successfully!**")

if uploaded_file is not None:
    e_name = uploaded_file.name
    file_extension = e_name.split(".")[-1].lower()
    text = ""
    if file_extension == "txt":
        text = uploaded_file.read().decode("utf-8")
    elif file_extension == "pdf":
        with pdfplumber.open(io.BytesIO(uploaded_file.read())) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    elif file_extension == "docx":
        doc = Document(io.BytesIO(uploaded_file.read()))
        text = "\n".join([p.text for p in doc.paragraphs])
    elif file_extension == "pptx":
        ppt = Presentation(io.BytesIO(uploaded_file.read()))
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n" 
    else:
        st.error("Unsupported file type!")

    # Only run if text was extracted
    if text:

        # ---- STATS ----
        words = text.split()
        sentences = re.split(r'[.!?]', text)
        sentences = [s for s in sentences if s.strip()]
        word_counts = Counter(w.lower().strip(".,!?") for w in words)

        # ---- DISPLAY ----
        st.markdown("<h3 style='text-align: center;'>📊 File Stats</h3>", unsafe_allow_html=True)

        col1, col2, col3 = st.columns(3)
        col1.metric("Words", len(words))
        col2.metric("Characters", len(text))
        col3.metric("Sentences", len(sentences))

        # ---- TOP WORDS CHART ----
        st.markdown("<h3 style='text-align: center;'>🔤 Top 10 Most Frequent Words</h3>", unsafe_allow_html=True)
        top_words = dict(word_counts.most_common(10))
        st.bar_chart(top_words)

        # ---- DOWNLOAD ----
        summary = f"""Text Inspector - Summary Report
================================
File Name  : {e_name}
Words      : {len(words)}
Characters : {len(text)}
Sentences  : {len(sentences)}

Top 10 Most Frequent Words:
{chr(10).join(f"{word}: {count}" for word, count in word_counts.most_common(10))}
"""

        st.download_button(
            label="⬇️ Download Summary",
            data=summary,
            file_name="text_inspector_report.txt",
            mime="text/plain"
        )

    else:
        st.warning("No text found in this file!")